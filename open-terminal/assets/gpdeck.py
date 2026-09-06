"""
gpdeck - branded presentation helper for the Open Terminal sandbox.

Turns a few lines of Python into a professional, on-brand .pptx (and PDF).
The AI agent should ALWAYS build decks through this module instead of writing
raw python-pptx from scratch (which produces ugly default-white slides).

Quick start
-----------
    import gpdeck
    d = gpdeck.Deck("pristatymas")
    d.title("Pavadinimas", "Paantraštė")
    d.bullets("Apie mus", ["Punktas vienas", "Punktas du", "Punktas trys"])
    img = d.gen_image("modern office illustration, clean, blue tones")
    d.image("Paslaugos", img, caption="Iliustracija")
    chart = d.chart_bar("Rodikliai", {"2022": 120, "2023": 145, "2024": 190})
    d.image("Augimas", chart)
    d.closing("Ačiū!", "www.pavyzdys.lt")
    d.save()          # -> pristatymas.pptx  (editable PowerPoint)
    d.to_pdf()        # -> pristatymas.pdf   (via LibreOffice)

Branding is read from brand.json / ~/.gpbrand/brand.json at runtime (never from
this source file) and can be overridden per deck:
gpdeck.Deck("name", brand={"primary": "1B4F9C", "company": "Acme"}).
"""

import json
import os
import subprocess
import urllib.request

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

_HERE = os.path.dirname(os.path.abspath(__file__))
PROXY = os.getenv("OT_PROXY_URL", "http://open-terminal-proxy:8000")

# Intermediate files (charts, generated images) default here — OUTSIDE the
# working dir — so only the final .pptx/.pdf appear in the user's file browser.
# They get embedded into the deck at save time; the temp copies don't matter.
WORK_DIR = os.getenv("GP_WORK_DIR", "/tmp/gpwork")
os.makedirs(WORK_DIR, exist_ok=True)

# "AI GENERATED" transparency label, stamped on every slide (compliance).
# dark = black label (for light slides), light = white label (for dark slides).
AI_LABEL_DARK = os.path.join(_HERE, "ai_label_dark.png")
AI_LABEL_LIGHT = os.path.join(_HERE, "ai_label_light.png")

# ---- branding -------------------------------------------------------------

# Generic defaults. The REAL per-deployment brand (company name, colours, logos)
# is loaded from brand.json / ~/.gpbrand/brand.json at runtime — it never lives
# in this (published) source file.
_DEFAULT_BRAND = {
    "company": "Company",
    "primary": "1F3C88",     # deep blue
    "secondary": "2E6BD6",   # bright blue
    "accent": "F5A623",      # amber accent
    "dark": "1A2233",        # near-black navy (text)
    "light": "F4F7FB",       # very light blue-grey (backgrounds)
    "muted": "6B7A90",       # grey for captions
    "font": "Inter",         # body font (installed in image)
    "font_head": "Inter",    # heading font
    "logo": "",              # logo png for LIGHT slides (empty = text logo)
    "logo_white": "",        # logo png for DARK slides
    # surface tokens (empty = derived from the palette above in _load_brand);
    # themes flip these to change the whole look without touching slide code
    "bg": "",                # content-slide background
    "text": "",              # body text on content slides
    "card": "",              # card / panel fill
    "heading": "",           # section-title colour on content slides
}

# Named visual themes the end user can pick ("padaryk tamsia tema" -> "dark").
# Each is a partial brand override applied on top of the deployment brand.
THEMES = {
    "corporate": {},  # the deployment brand as-is (default)
    "dark": {
        "bg": "0F1A2B", "text": "E7EDF4", "card": "1C2B42", "heading": "FFFFFF",
        "muted": "8CA0B6", "primary": "2E6BD6", "secondary": "4C8DE0",
        "light": "1C2B42",
    },
    "bold": {
        "primary": "111827", "secondary": "374151", "accent": "FB7185",
        "heading": "111827", "dark": "111827",
    },
    "minimal": {
        "primary": "1A1A1A", "secondary": "4A4A4A", "accent": "C69214",
        "muted": "9AA0A6", "card": "FAFAFA", "heading": "1A1A1A",
        "dark": "222222",
    },
    "emerald": {
        "primary": "0F5132", "secondary": "198754", "accent": "E9B949",
        "heading": "0F5132",
    },
}

# Where a deployment's brand.json + logo files may live. First existing wins as
# the override on top of the generic default; logo paths are resolved relative to
# whichever file supplied them, so brand + logos travel together in the volume.
# NOTE: the container runs as root, so ~ expands to /root, but the sandbox's
# user-facing home (file-browser root, where work happens) is /home/user — so we
# look there explicitly rather than relying on $HOME.
_BRAND_PATHS = [
    os.path.join(_HERE, "brand.json"),                       # baked next to gpdeck
    "/home/user/.gpbrand/brand.json",                        # sandbox volume (home)
    os.path.expanduser("~/.gpbrand/brand.json"),             # $HOME fallback
    os.path.join(os.getcwd(), ".gpbrand", "brand.json"),     # cwd fallback
    os.path.expanduser("~/.gpbrand.json"),                   # legacy single file
]


def _load_brand(override=None, theme=None):
    brand = dict(_DEFAULT_BRAND)
    # start from a generic default file if present (published template)
    for base in (os.path.join(_HERE, "brand.default.json"),):
        try:
            with open(base, "r", encoding="utf-8") as f:
                brand.update({k: v for k, v in json.load(f).items()
                              if v and not k.startswith("_")})
        except Exception:
            pass
    # then the real per-deployment brand, resolving logo paths next to it
    for path in _BRAND_PATHS:
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception:
            continue
        d = os.path.dirname(os.path.abspath(path))
        for key in ("logo", "logo_white"):
            v = data.get(key)
            if v and not os.path.isabs(v):
                data[key] = os.path.join(d, v)
        brand.update({k: v for k, v in data.items()
                      if v and not k.startswith("_")})
        break
    # named theme on top of the deployment brand
    if theme:
        brand.update({k: v for k, v in THEMES.get(theme, {}).items() if v})
    if override:
        brand.update({k: v for k, v in override.items() if v})
    # derive surface tokens from the palette when a theme didn't set them
    brand["bg"] = brand.get("bg") or "FFFFFF"
    brand["text"] = brand.get("text") or brand["dark"]
    brand["card"] = brand.get("card") or brand["light"]
    brand["heading"] = brand.get("heading") or brand["primary"]
    return brand


def _rgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#"))


def _is_dark(hexstr):
    """True if a hex colour is dark (so labels/logos should go light)."""
    h = hexstr.lstrip("#")[:6]
    try:
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return False
    return (0.299 * r + 0.587 * g + 0.114 * b) < 128


def gen_image(prompt, filename=None, model=None):
    """Generate an image via the Gemini image model (proxy holds the API key)
    and save it to the working dir. Returns the file path, or None on failure."""
    filename = filename or os.path.join(WORK_DIR, f"gen_{abs(hash(prompt)) % 10**8}.png")
    if not filename.lower().endswith((".png", ".jpg", ".jpeg")):
        filename += ".png"
    body = json.dumps({"prompt": prompt, "model": model} if model
                      else {"prompt": prompt}).encode()
    req = urllib.request.Request(PROXY + "/genimage", data=body,
                                 headers={"Content-Type": "application/json"})
    try:
        import base64
        with urllib.request.urlopen(req, timeout=180) as r:
            d = json.load(r)
        raw = base64.b64decode(d["b64"])
        with open(filename, "wb") as f:
            f.write(raw)
        return filename
    except Exception as e:
        print(f"[gpdeck] gen_image failed: {e}")
        return None


# ---- brand kit (official assets synced from SharePoint) -------------------

KIT_DIR = os.getenv("GP_KIT_DIR", "/home/user/.gpbrand/kit")
_IMG_EXTS = (".png", ".jpg", ".jpeg", ".gif", ".webp")


def kit_images(contains=None):
    """List official brand images synced into the kit (absolute paths).
    Optional `contains` filters by substring of the path (case-insensitive)."""
    out = []
    manifest = os.path.join(KIT_DIR, "manifest.json")
    if os.path.exists(manifest):
        try:
            with open(manifest, encoding="utf-8") as f:
                for rel in json.load(f).get("images", []):
                    out.append(os.path.join(KIT_DIR, rel))
        except Exception:
            pass
    if not out and os.path.isdir(KIT_DIR):
        for root, _d, files in os.walk(KIT_DIR):
            for fn in files:
                if fn.lower().endswith(_IMG_EXTS):
                    out.append(os.path.join(root, fn))
    if contains:
        c = contains.lower()
        out = [p for p in out if c in p.lower()]
    return sorted(out)


def kit_templates():
    """List official PowerPoint templates (.pptx/.potx) in the kit."""
    out = []
    if os.path.isdir(KIT_DIR):
        for root, _d, files in os.walk(KIT_DIR):
            for fn in files:
                if fn.lower().endswith((".pptx", ".potx")):
                    out.append(os.path.join(root, fn))
    return sorted(out)


# ---- icons (lucide SVG, recoloured to brand + rasterised to PNG) ----------

ICON_DIR = os.getenv("GP_ICON_DIR", "/opt/gpdeck/icons")
_ICON_CACHE = "/tmp/gpicons"


def icons_available():
    """Names of bundled icons (call to discover valid icon names)."""
    if not os.path.isdir(ICON_DIR):
        return []
    return sorted(f[:-4] for f in os.listdir(ICON_DIR) if f.endswith(".svg"))


def icon(name, color="1F3C88", size=128):
    """Return a PNG path for icon `name` recoloured to `color` (hex), or None.
    Icons are lucide line icons; `gpdeck.icons_available()` lists valid names."""
    if not name:
        return None
    src = os.path.join(ICON_DIR, name + ".svg")
    if not os.path.exists(src):
        return None
    color = "#" + color.lstrip("#")
    os.makedirs(_ICON_CACHE, exist_ok=True)
    out = os.path.join(_ICON_CACHE, f"{name}_{color.lstrip('#')}_{size}.png")
    if os.path.exists(out):
        return out
    try:
        import cairosvg
        with open(src, encoding="utf-8") as f:
            svg = f.read()
        svg = svg.replace("currentColor", color)
        cairosvg.svg2png(bytestring=svg.encode(), write_to=out,
                         output_width=size, output_height=size)
        return out
    except Exception:
        return None


# 16:9 canvas
SW = Inches(13.333)
SH = Inches(7.5)


class Deck:
    def __init__(self, name="presentation", brand=None, theme=None, template=None):
        """theme: one of gpdeck.THEMES (corporate/dark/bold/minimal/emerald).
        template: path to an official .pptx/.potx to build on (inherits its
        masters/theme); use gpdeck.kit_templates() to list synced ones."""
        self.name = name[:-5] if name.lower().endswith(".pptx") else name
        self.b = _load_brand(brand, theme=theme)
        if template and os.path.exists(template):
            self.prs = Presentation(template)
        else:
            self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        try:
            self._blank = self.prs.slide_layouts[6]     # standard "Blank" layout
        except IndexError:
            self._blank = self.prs.slide_layouts[-1]
        self.theme = theme
        self._pageno = 0          # content-slide counter for the footer
        self.footer = True        # set False to hide footers/page numbers
        self.ai_label = True      # stamp the "AI GENERATED" label on every slide

    # ---- low-level helpers ----------------------------------------------

    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _rect(self, slide, x, y, w, h, color, line=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(color)
        if line:
            shp.line.color.rgb = _rgb(line)
            shp.line.width = Pt(1)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(self, slide, x, y, w, h, text, size=18, color=None, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None,
              line_spacing=1.15):
        color = color or self.b["text"]
        font = font or self.b["font"]
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = 0
        tf.margin_right = 0
        lines = text if isinstance(text, list) else [text]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = str(ln)
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = _rgb(color)
        return tb

    def _bg(self, slide, color):
        self._rect(slide, 0, 0, SW, SH, color)

    def _logo(self, slide, x, y, h_in=0.42, on_dark=False):
        logo = self.b.get("logo_white") if on_dark else self.b.get("logo")
        logo = logo or self.b.get("logo")  # fall back to the single logo
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, x, y, height=Inches(h_in))
                return
            except Exception:
                pass
        # text logo fallback
        self._text(slide, x, y, Inches(4), Inches(0.5), self.b["company"],
                   size=14, bold=True,
                   color="FFFFFF" if on_dark else self.b["primary"],
                   font=self.b["font_head"])

    # ---- slide types -----------------------------------------------------

    def title(self, title, subtitle=""):
        s = self._slide()
        self._bg(s, self.b["primary"])
        # accent bar
        self._rect(s, 0, Inches(4.55), Inches(2.2), Inches(0.09), self.b["accent"])
        self._text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.8),
                   title, size=48, bold=True, color="FFFFFF",
                   font=self.b["font_head"])
        if subtitle:
            self._text(s, Inches(0.92), Inches(4.75), Inches(11.3), Inches(1.2),
                       subtitle, size=22, color="D8E3F3", font=self.b["font"])
        self._logo(s, Inches(0.9), Inches(0.7), on_dark=True)
        self._ai(s, on_dark=True)
        return s

    def section(self, title, number=None):
        s = self._slide()
        self._bg(s, self.b["light"])
        self._rect(s, 0, Inches(3.15), SW, Inches(1.5), self.b["secondary"])
        if number:
            self._text(s, Inches(0.9), Inches(2.2), Inches(3), Inches(1),
                       str(number), size=64, bold=True, color=self.b["accent"],
                       font=self.b["font_head"])
        self._text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(1.1),
                   title, size=34, bold=True, color="FFFFFF",
                   anchor=MSO_ANCHOR.MIDDLE, font=self.b["font_head"])
        self._ai(s, on_dark=_is_dark(self.b["light"]))
        return s

    def _header(self, s):
        """Content slide with a branded title band; returns content-top (Emu)."""
        self._bg(s, self.b["bg"])
        self._rect(s, 0, 0, Inches(0.28), SH, self.b["primary"])  # left spine
        return s

    def _title_band(self, s, title):
        self._text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.0),
                   title, size=30, bold=True, color=self.b["heading"],
                   font=self.b["font_head"])
        self._rect(s, Inches(0.92), Inches(1.5), Inches(1.6), Inches(0.06),
                   self.b["accent"])

    def bullets(self, title, items):
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        tb = slide_tb = self._text(
            s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.9),
            "", size=20, color=self.b["text"])
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing = 1.2
            p.space_after = Pt(14)
            # bullet glyph in accent
            rb = p.add_run(); rb.text = "•  "
            rb.font.size = Pt(20); rb.font.bold = True
            rb.font.color.rgb = _rgb(self.b["accent"]); rb.font.name = self.b["font"]
            rt = p.add_run(); rt.text = str(it)
            rt.font.size = Pt(20); rt.font.color.rgb = _rgb(self.b["text"])
            rt.font.name = self.b["font"]
        self._footer(s)
        return s

    def two_column(self, title, left_title, left_items, right_title, right_items):
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        for x, ct, items in ((Inches(1.0), left_title, left_items),
                             (Inches(7.0), right_title, right_items)):
            self._rect(s, x, Inches(2.05), Inches(5.3), Inches(0.55),
                       self.b["card"])
            self._text(s, x + Inches(0.2), Inches(2.1), Inches(5.0), Inches(0.5),
                       ct, size=18, bold=True, color=self.b["heading"],
                       anchor=MSO_ANCHOR.MIDDLE, font=self.b["font_head"])
            tb = self._text(s, x + Inches(0.2), Inches(2.8), Inches(5.0),
                            Inches(4.0), "", size=17)
            tf = tb.text_frame; tf.word_wrap = True
            for i, it in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10); p.line_spacing = 1.15
                rb = p.add_run(); rb.text = "•  "
                rb.font.size = Pt(17); rb.font.color.rgb = _rgb(self.b["accent"])
                rt = p.add_run(); rt.text = str(it)
                rt.font.size = Pt(17); rt.font.color.rgb = _rgb(self.b["text"])
                rt.font.name = self.b["font"]
        self._footer(s)
        return s

    def image(self, title, image_path, caption=""):
        """Title slide with a large image on the right / centered."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        if image_path and os.path.exists(image_path):
            # fit into a 11.3 x 4.6 box, centered, keep aspect
            self._place_image(s, image_path, Inches(1.0), Inches(2.0),
                              Inches(11.3), Inches(4.5))
        if caption:
            self._text(s, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.4),
                       caption, size=13, color=self.b["muted"],
                       align=PP_ALIGN.CENTER)
        self._footer(s)
        return s

    def image_full(self, image_path, title="", subtitle=""):
        """Full-bleed image with an optional dark gradient title bar."""
        s = self._slide()
        self._bg(s, "1A2233")   # neutral dark behind the full-bleed image
        if image_path and os.path.exists(image_path):
            self._place_image(s, image_path, 0, 0, SW, SH, cover=True)
        if title:
            self._rect(s, 0, Inches(5.6), SW, Inches(1.9), self.b["primary"])
            # simulate translucency by a solid band
            self._text(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.0),
                       title, size=32, bold=True, color="FFFFFF",
                       font=self.b["font_head"])
            if subtitle:
                self._text(s, Inches(0.92), Inches(6.7), Inches(11.3),
                           Inches(0.6), subtitle, size=16, color="D8E3F3")
        self._ai(s, on_dark=True)
        return s

    def quote(self, text, author=""):
        s = self._slide()
        self._bg(s, self.b["light"])
        self._text(s, Inches(1.2), Inches(1.0), Inches(3), Inches(2), "“",
                   size=120, bold=True, color=self.b["accent"],
                   font=self.b["font_head"])
        self._text(s, Inches(1.4), Inches(2.6), Inches(10.5), Inches(2.8),
                   text, size=28, bold=True, color=self.b["primary"],
                   font=self.b["font_head"])
        if author:
            self._text(s, Inches(1.4), Inches(5.6), Inches(10.5), Inches(0.6),
                       "— " + author, size=18, color=self.b["muted"])
        self._ai(s, on_dark=_is_dark(self.b["light"]))
        return s

    def closing(self, title="Ačiū!", subtitle=""):
        s = self._slide()
        self._bg(s, self.b["primary"])
        self._rect(s, Inches(0.9), Inches(3.45), Inches(2.2), Inches(0.09),
                   self.b["accent"])
        self._text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.4),
                   title, size=44, bold=True, color="FFFFFF",
                   font=self.b["font_head"])
        if subtitle:
            self._text(s, Inches(0.92), Inches(3.7), Inches(11.3), Inches(1.0),
                       subtitle, size=20, color="D8E3F3")
        self._logo(s, Inches(0.9), Inches(6.4), on_dark=True)
        self._ai(s, on_dark=True)
        return s

    # ---- "AI GENERATED" transparency label -------------------------------

    def _ai(self, slide, on_dark=False):
        """Stamp the AI-GENERATED label top-right of a slide."""
        if not self.ai_label:
            return
        path = AI_LABEL_LIGHT if on_dark else AI_LABEL_DARK
        if not os.path.exists(path):
            return
        h = Inches(0.22)
        w = Inches(0.22 * 3.156)   # label aspect ratio 7459x2363
        try:
            slide.shapes.add_picture(path, int(SW - w - Inches(0.35)),
                                     Inches(0.3), width=int(w), height=int(h))
        except Exception:
            pass

    # ---- footer / page numbers (content slides) --------------------------

    def _footer(self, slide):
        # every content slide gets the footer AND the AI label
        dark = _is_dark(self.b["bg"])
        self._ai(slide, on_dark=dark)
        if not self.footer:
            return
        self._pageno += 1
        y = Inches(7.02)
        logo = (self.b.get("logo_white") if dark else self.b.get("logo")) \
            or self.b.get("logo")
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(0.9), y, height=Inches(0.22))
            except Exception:
                pass
        else:
            self._text(slide, Inches(0.9), Inches(6.98), Inches(4), Inches(0.4),
                       self.b["company"], size=9, bold=True,
                       color=self.b["muted"], font=self.b["font_head"])
        self._text(slide, Inches(9.3), Inches(6.98), Inches(3.13), Inches(0.4),
                   f"{self.b['company']} · {self._pageno}", size=9,
                   color=self.b["muted"], align=PP_ALIGN.RIGHT)

    # ---- more slide types ------------------------------------------------

    def kpi(self, title, stats):
        """Big-number highlights. stats = [(value, label), ...] (2-4 items)."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        stats = list(stats)[:4]
        n = len(stats) or 1
        gap = Inches(0.4)
        total_w = Inches(11.3)
        cw = (total_w - gap * (n - 1)) / n
        x = Inches(1.0)
        for value, label in stats:
            self._rect(s, x, Inches(2.5), cw, Inches(2.9), self.b["card"])
            self._rect(s, x, Inches(2.5), cw, Inches(0.12), self.b["accent"])
            self._text(s, x, Inches(3.0), cw, Inches(1.3), str(value), size=54,
                       bold=True, color=self.b["heading"], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE, font=self.b["font_head"])
            self._text(s, x + Inches(0.2), Inches(4.35), cw - Inches(0.4),
                       Inches(0.9), str(label), size=15, color=self.b["text"],
                       align=PP_ALIGN.CENTER)
            x = int(x + cw + gap)
        self._footer(s)
        return s

    def agenda(self, title, items):
        """Numbered table of contents."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        y = Inches(2.1)
        for i, it in enumerate(items[:7], 1):
            self._rect(s, Inches(1.0), y, Inches(0.55), Inches(0.55),
                       self.b["primary"])
            self._text(s, Inches(1.0), y, Inches(0.55), Inches(0.55), str(i),
                       size=20, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE, font=self.b["font_head"])
            self._text(s, Inches(1.8), y, Inches(10.4), Inches(0.55), str(it),
                       size=20, color=self.b["text"], anchor=MSO_ANCHOR.MIDDLE)
            y = int(y + Inches(0.72))
        self._footer(s)
        return s

    def comparison(self, title, left_title, left_items, right_title,
                   right_items, left_good=True):
        """Two columns with check / cross marks (e.g. pros vs cons)."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        cols = ((Inches(1.0), left_title, left_items, self.b["secondary"],
                 "✓" if left_good else "✗"),
                (Inches(7.0), right_title, right_items, self.b["muted"],
                 "✗" if left_good else "✓"))
        for x, ct, items, color, mark in cols:
            self._rect(s, x, Inches(2.05), Inches(5.3), Inches(0.6), color)
            self._text(s, x, Inches(2.1), Inches(5.3), Inches(0.5), ct, size=18,
                       bold=True, color="FFFFFF", align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE, font=self.b["font_head"])
            tb = self._text(s, x + Inches(0.3), Inches(2.9), Inches(4.8),
                            Inches(4.0), "", size=17)
            tf = tb.text_frame; tf.word_wrap = True
            for i, it in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10); p.line_spacing = 1.15
                rm = p.add_run(); rm.text = mark + "  "
                rm.font.size = Pt(17); rm.font.bold = True
                rm.font.color.rgb = _rgb(color); rm.font.name = self.b["font"]
                rt = p.add_run(); rt.text = str(it)
                rt.font.size = Pt(17); rt.font.color.rgb = _rgb(self.b["text"])
                rt.font.name = self.b["font"]
        self._footer(s)
        return s

    def timeline(self, title, steps):
        """Horizontal timeline. steps = [(label, text), ...] (2-5)."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        steps = list(steps)[:5]
        n = len(steps) or 1
        self._rect(s, Inches(1.2), Inches(3.55), Inches(10.9), Inches(0.05),
                   self.b["secondary"])
        span = Inches(10.9) / max(n - 1, 1) if n > 1 else 0
        x0 = Inches(1.2)
        for i, (label, text) in enumerate(steps):
            cx = int(x0 + span * i) if n > 1 else int(x0 + Inches(5.45))
            self._rect(s, cx - Inches(0.16), Inches(3.4), Inches(0.34),
                       Inches(0.34), self.b["accent"])
            top = (i % 2 == 0)
            ly = Inches(2.1) if top else Inches(4.0)
            self._text(s, cx - Inches(1.3), ly, Inches(2.6), Inches(0.4),
                       str(label), size=16, bold=True, color=self.b["primary"],
                       align=PP_ALIGN.CENTER, font=self.b["font_head"])
            self._text(s, cx - Inches(1.3), int(ly + Inches(0.4)), Inches(2.6),
                       Inches(1.1), str(text), size=12, color=self.b["text"],
                       align=PP_ALIGN.CENTER)
        self._footer(s)
        return s

    def table(self, title, headers, rows):
        """Branded table. headers = [..], rows = [[..], ..]."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        ncol = len(headers)
        nrow = len(rows) + 1
        gt = slide_tbl = s.shapes.add_table(
            nrow, ncol, Inches(1.0), Inches(2.1), Inches(11.3),
            Inches(min(4.6, 0.5 * nrow))).table
        for c, h in enumerate(headers):
            cell = gt.cell(0, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(self.b["primary"])
            p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = str(h)
            r.font.size = Pt(14); r.font.bold = True
            r.font.color.rgb = _rgb("FFFFFF"); r.font.name = self.b["font_head"]
        for ri, row in enumerate(rows, 1):
            for c in range(ncol):
                cell = gt.cell(ri, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(
                    self.b["bg"] if ri % 2 else self.b["card"])
                val = row[c] if c < len(row) else ""
                p = cell.text_frame.paragraphs[0]; r = p.add_run()
                r.text = str(val)
                r.font.size = Pt(13); r.font.color.rgb = _rgb(self.b["text"])
                r.font.name = self.b["font"]
        self._footer(s)
        return s

    def icon_grid(self, title, items):
        """Grid of icon + label + text. items = [(icon_name, label, text), ...].
        Icons come from gpdeck.icon(); missing icons degrade to an accent dot."""
        s = self._slide()
        self._header(s)
        self._title_band(s, title)
        items = list(items)[:6]
        cols = 3 if len(items) > 2 else len(items) or 1
        cw = Inches(3.7); ch = Inches(2.15)
        x0, y0 = Inches(1.0), Inches(2.1)
        for i, item in enumerate(items):
            name, label, text = (list(item) + ["", "", ""])[:3]
            r, c = divmod(i, cols)
            x = int(x0 + c * (cw + Inches(0.1)))
            y = int(y0 + r * (ch + Inches(0.2)))
            ipath = icon(name, self.b["primary"]) if name else None
            if ipath and os.path.exists(ipath):
                try:
                    s.shapes.add_picture(ipath, x, y, height=Inches(0.6))
                except Exception:
                    ipath = None
            if not ipath:
                self._rect(s, x, y, Inches(0.35), Inches(0.35), self.b["accent"])
            self._text(s, x, int(y + Inches(0.72)), cw - Inches(0.2), Inches(0.5),
                       str(label), size=17, bold=True, color=self.b["primary"],
                       font=self.b["font_head"])
            self._text(s, x, int(y + Inches(1.2)), cw - Inches(0.2), Inches(0.9),
                       str(text), size=12, color=self.b["text"])
        self._footer(s)
        return s

    # ---- image placement -------------------------------------------------

    def _place_image(self, slide, path, x, y, w, h, cover=False):
        from PIL import Image
        try:
            iw, ih = Image.open(path).size
        except Exception:
            slide.shapes.add_picture(path, x, y, width=w)
            return
        box_ratio = w / h
        img_ratio = iw / ih
        if cover:
            # fill the box, crop overflow
            if img_ratio > box_ratio:
                new_h = h; new_w = int(h * img_ratio)
            else:
                new_w = w; new_h = int(w / img_ratio)
            px = x + (w - new_w) / 2
            py = y + (h - new_h) / 2
            slide.shapes.add_picture(path, int(px), int(py), int(new_w), int(new_h))
        else:
            # fit inside the box, centered
            if img_ratio > box_ratio:
                new_w = w; new_h = int(w / img_ratio)
            else:
                new_h = h; new_w = int(h * img_ratio)
            px = x + (w - new_w) / 2
            py = y + (h - new_h) / 2
            slide.shapes.add_picture(path, int(px), int(py), int(new_w), int(new_h))

    # ---- AI image generation --------------------------------------------

    def gen_image(self, prompt, filename=None, model=None):
        """Generate an image with the Gemini image model (via the proxy) and
        save it into the working dir. Returns the file path (or None). The API
        key never touches the sandbox — the proxy holds it."""
        return gen_image(prompt, filename, model)

    # ---- charts (branded matplotlib) ------------------------------------

    def _mpl(self):
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        plt.rcParams.update({
            "font.size": 13,
            "axes.edgecolor": "#" + self.b["muted"],
            "axes.labelcolor": "#" + self.b["text"],
            "text.color": "#" + self.b["text"],
            "xtick.color": "#" + self.b["text"],
            "ytick.color": "#" + self.b["text"],
            "axes.grid": True,
            "grid.color": "#" + self.b["muted"] + "44",
            "figure.facecolor": "#" + self.b["bg"],
            "axes.facecolor": "#" + self.b["bg"],
        })
        return plt

    def _palette(self):
        return ["#" + self.b[k] for k in ("secondary", "accent", "primary",
                                          "muted")] + ["#7FB3E8", "#F5C77E"]

    def chart_bar(self, title, data, filename=None, ylabel=""):
        plt = self._mpl()
        fig, ax = plt.subplots(figsize=(9, 5), dpi=200)
        keys = list(data.keys()); vals = list(data.values())
        ax.bar(keys, vals, color="#" + self.b["secondary"], width=0.6)
        ax.set_title(title, fontsize=16, fontweight="bold", color="#" + self.b["primary"])
        if ylabel:
            ax.set_ylabel(ylabel)
        ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        return self._savefig(plt, fig, filename, "bar")

    def chart_line(self, title, series, filename=None, ylabel=""):
        """series: {'label': {x: y, ...}, ...} or {x: y} for a single line."""
        plt = self._mpl()
        fig, ax = plt.subplots(figsize=(9, 5), dpi=200)
        pal = self._palette()
        if series and isinstance(next(iter(series.values())), dict):
            for i, (label, d) in enumerate(series.items()):
                ax.plot(list(d.keys()), list(d.values()), marker="o",
                        linewidth=2.5, color=pal[i % len(pal)], label=label)
            ax.legend(frameon=False)
        else:
            ax.plot(list(series.keys()), list(series.values()), marker="o",
                    linewidth=2.5, color="#" + self.b["secondary"])
        ax.set_title(title, fontsize=16, fontweight="bold", color="#" + self.b["primary"])
        if ylabel:
            ax.set_ylabel(ylabel)
        ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        return self._savefig(plt, fig, filename, "line")

    def chart_pie(self, title, data, filename=None):
        plt = self._mpl()
        fig, ax = plt.subplots(figsize=(7, 6), dpi=200)
        ax.pie(list(data.values()), labels=list(data.keys()), autopct="%1.0f%%",
               colors=self._palette(), startangle=90,
               wedgeprops={"edgecolor": "white", "linewidth": 2})
        ax.set_title(title, fontsize=16, fontweight="bold", color="#" + self.b["primary"])
        return self._savefig(plt, fig, filename, "pie")

    def _savefig(self, plt, fig, filename, kind):
        filename = filename or os.path.join(
            WORK_DIR, f"chart_{kind}_{id(fig) % 10**6}.png")
        if not filename.lower().endswith(".png"):
            filename += ".png"
        fig.tight_layout()
        fig.savefig(filename, bbox_inches="tight", facecolor="#" + self.b["bg"])
        plt.close(fig)
        return filename

    # ---- output ----------------------------------------------------------

    def save(self, path=None):
        path = path or (self.name + ".pptx")
        if not path.lower().endswith(".pptx"):
            path += ".pptx"
        self.prs.save(path)
        print(f"[gpdeck] saved {path} ({len(self.prs.slides._sldIdLst)} slides)")
        return path

    def to_pdf(self, pptx_path=None, outdir="."):
        pptx_path = pptx_path or (self.name + ".pptx")
        if not os.path.exists(pptx_path):
            self.save(pptx_path)
        r = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir",
             outdir, pptx_path],
            capture_output=True, text=True, timeout=180)
        pdf = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
        pdf = os.path.join(outdir, pdf)
        if os.path.exists(pdf):
            print(f"[gpdeck] pdf {pdf}")
            return pdf
        print(f"[gpdeck] pdf failed: {r.stderr[:300]}")
        return None
